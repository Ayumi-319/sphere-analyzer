import streamlit as st
from cellpose import models
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from skimage import measure, segmentation, exposure
import io as python_io
from pptx import Presentation
from PIL import Image
import torch
import cv2

st.set_page_config(page_title="Sphere Analyzer v2.5.1", layout="wide")
st.title("🔴 スフェア自動計測ツール v2.5.1")

@st.cache_resource
def get_model():
    return models.CellposeModel(gpu=False, model_type='cyto', device=torch.device('cpu'))

with st.sidebar:
    st.header("1. 基本設定")
    mag = st.radio("倍率:", ("4x", "10x", "カスタム"), index=0)
    um_per_pixel = 3.23 if mag == "4x" else 1.28 if mag == "10x" else st.number_input("μm/px", value=1.0)

    st.header("2. 画像補正 (再計算)")
    # コントラスト調整のアルゴリズムを変更
    contrast_factor = st.slider("コントラスト倍率", 0.0, 2.0, 1.0, help="1.0より上げると暗い部分がより暗く、エッジが際立ちます")
    
    st.header("3. AI解析設定 (再計算)")
    target_diameter = st.number_input("予想直径 (px)", value=150)
    flow_threshold = st.slider("切り離し強度", 0.0, 1.1, 0.4)
    
    st.header("4. フィルタ設定 (即時反映)")
    exclude_border = st.checkbox("画像端の個体を除外", value=True)
    circularity_threshold = st.slider("真円度しきい値", 0.0, 1.0, 0.8)

uploaded_files = st.file_uploader("ファイルをドロップ", type=['jpg', 'png', 'pptx', 'jpeg'], accept_multiple_files=True)

if 'analysis_cache' not in st.session_state:
    st.session_state.analysis_cache = {}

if uploaded_files:
    model = get_model()
    all_final_results = []

    for f in uploaded_files:
        if f.name.lower().endswith('.pptx'):
            prs = Presentation(f)
            images = [(f"{f.name}_S{i+1}", np.array(Image.open(python_io.BytesIO(s.image.blob)))) 
                      for i, sl in enumerate(prs.slides) for s in sl.shapes if s.shape_type == 13]
        else:
            images = [(f.name, np.array(Image.open(f)))]

        for name, original_img in images:
            st.subheader(f"解析: {name}")
            
            cache_key = f"{name}_{contrast_factor}_{target_diameter}_{flow_threshold}"
            if cache_key not in st.session_state.analysis_cache:
                with st.spinner(f'{name} を解析中...'):
                    # 画像補正：コントラスト強調
                    # 暗い部分をカットしてエッジを立たせる処理
                    p2, p98 = np.percentile(original_img, (2 * contrast_factor, 98))
                    img_adj = exposure.rescale_intensity(original_img, in_range=(p2, p98))
                    
                    h, w = img_adj.shape[:2]
                    scale = 800 / max(h, w) if max(h, w) > 800 else 1.0
                    resized_img = cv2.resize(img_adj, (int(w*scale), int(h*scale)))
                    
                    masks, _, _ = model.eval(resized_img, 
                                             diameter=target_diameter*scale, 
                                             flow_threshold=flow_threshold,
                                             channels=[0,0])
                    
                    masks = cv2.resize(masks.astype(np.uint16), (w, h), interpolation=cv2.INTER_NEAREST)
                    
                    if exclude_border:
                        masks = segmentation.clear_border(masks)
                        
                    props = measure.regionprops_table(masks, properties=['label', 'area', 'perimeter', 'equivalent_diameter'])
                    st.session_state.analysis_cache[cache_key] = (pd.DataFrame(props), masks)

            full_df, masks = st.session_state.analysis_cache[cache_key]
            
            df = full_df.copy()
            if not df.empty:
                df['circularity'] = (4 * np.pi * df['area']) / (df['perimeter'] ** 2)
                df['diameter_um'] = df['equivalent_diameter'] * um_per_pixel
                df['filename'] = name
                df_clean = df[df['circularity'] > circularity_threshold].copy()
                all_final_results.append(df_clean)

                col1, col2 = st.columns([2, 1])
                with col1:
                    fig, ax = plt.subplots()
                    ax.imshow(original_img)
                    ax.contour(masks > 0, colors='lime', linewidths=0.5)
                    ax.axis('off')
                    st.pyplot(fig)
                    plt.close(fig)
                with col2:
                    st.metric("検出数", f"{len(df_clean)} 個")
                    st.metric("平均直径", f"{df_clean['diameter_um'].mean():.1f} μm")
                    st.metric("平均真円度", f"{df_clean['circularity'].mean():.2f}")
            else:
                st.warning("スフェアが検出されませんでした。設定を調整してください。")

    if all_final_results:
        final_df = pd.concat(all_final_results)
        output = python_io.BytesIO()
        with pd.ExcelWriter(output, engine='openpyxl') as writer:
            final_df.to_excel(writer, index=False)
        st.download_button("📊 Excel保存", output.getvalue(), "results.xlsx")
